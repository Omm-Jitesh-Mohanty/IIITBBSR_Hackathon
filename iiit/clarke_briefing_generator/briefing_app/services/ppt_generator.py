
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import re
import requests
import os
from concurrent.futures import ThreadPoolExecutor


image_cache = {}


def download_image(query):

    if query in image_cache:
        return image_cache[query]

    try:

        seed = abs(hash(query)) % 1000

        # smaller image for faster download
        url = f"https://picsum.photos/seed/{seed}/640/480"

        response = requests.get(url, timeout=10)

        if response.status_code != 200:
            return None

        image_path = f"temp_{seed}.jpg"

        with open(image_path, "wb") as f:
            f.write(response.content)

        image_cache[query] = image_path

        return image_path

    except Exception as e:
        print("Image download failed:", e)
        return None


def create_ppt(ai_text):

    prs = Presentation()

    slides = re.split(r"\*\*Slide \d+: ", ai_text)[1:]

    

    image_queries = []

    parsed_slides = []

    for slide in slides:

        slide = slide.strip()
        lines = slide.split("\n")

        title = lines[0].replace("**", "").strip()

        bullets = []
        for line in lines[1:]:
            line = line.replace("-", "").strip()
            if line:
                bullets.append(line)

        bullets = bullets[:5]

        use_image = 2 <= len(bullets) <= 4

        query = None
        if use_image:
            query = " ".join(bullets[:2]) if bullets else title

        image_queries.append(query)

        parsed_slides.append((title, bullets, use_image))

    

    with ThreadPoolExecutor(max_workers=5) as executor:
        image_paths = list(
            executor.map(
                lambda q: download_image(q) if q else None,
                image_queries
            )
        )

    

    for i, slide_data in enumerate(parsed_slides):

        title, bullets, use_image = slide_data

        slide_obj = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        background = slide_obj.shapes.add_shape(
            1,
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(7.5)
        )

        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(20, 20, 20)
        background.line.fill.background()

        slide_obj.shapes._spTree.remove(background._element)
        slide_obj.shapes._spTree.insert(2, background._element)

        # Title
        title_box = slide_obj.shapes.add_textbox(
            Inches(0.7),
            Inches(0.4),
            Inches(12),
            Inches(1)
        )

        tf = title_box.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 60, 90)

        # Decorative bar
        bar = slide_obj.shapes.add_shape(
            1,
            Inches(0.7),
            Inches(1.25),
            Inches(3),
            Inches(0.15)
        )

        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(255, 60, 90)
        bar.line.fill.background()

        text_width = 6 if use_image else 11

        # Bullet box
        bullet_box = slide_obj.shapes.add_textbox(
            Inches(0.7),
            Inches(1.7),
            Inches(text_width),
            Inches(4.5)
        )

        tf = bullet_box.text_frame
        tf.word_wrap = True
        tf.clear()

        first = True

        for bullet in bullets:

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            if len(bullet) > 100:
                bullet = bullet[:100] + "..."

            p.text = "• " + bullet
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Add image (already downloaded)
        if use_image:

            img_path = image_paths[i]

            if img_path:

                try:
                    slide_obj.shapes.add_picture(
                        img_path,
                        Inches(7),
                        Inches(1.8),
                        width=Inches(5),
                        height=Inches(3.5)
                    )

                except Exception as e:
                    print("Image load failed:", e)

    file_path = "media/generated_presentation.pptx"

    prs.save(file_path)

    return file_path